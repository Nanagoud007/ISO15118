#!/usr/bin/env python3
"""
ISO 15118 vs OCPP 1.6 Presentation Generator
Creates a beautiful PowerPoint presentation from markdown content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import re

def create_iso_15118_presentation():
    """Create the ISO 15118 presentation"""
    
    # Create presentation
    prs = Presentation()
    
    # Define colors
    PRIMARY_BLUE = RGBColor(0, 102, 204)
    SECONDARY_GREEN = RGBColor(0, 204, 102)
    ACCENT_ORANGE = RGBColor(255, 102, 0)
    BACKGROUND_GRAY = RGBColor(245, 245, 245)
    TEXT_DARK = RGBColor(51, 51, 51)
    LIGHT_BLUE = RGBColor(173, 216, 230)
    LIGHT_GREEN = RGBColor(144, 238, 144)
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "ISO 15118: The Next Generation of EV Charging Communication"
    subtitle.text = "Comparing ISO 15118 with OCPP 1.6: Benefits for Vendors, Grid, and Users"
    
    # Apply formatting
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
    
    # Slide 2: Agenda
    slide_layout = prs.slide_layouts[1]  # Title and content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Agenda"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    agenda_items = [
        "Introduction to EV Charging Protocols",
        "OCPP 1.6 Overview & Architecture",
        "ISO 15118 Overview & Architecture",
        "System Architecture Comparison",
        "Workflow Diagrams",
        "Key Differences Comparison",
        "Advantages for Different Stakeholders",
        "Implementation Challenges",
        "Future Outlook",
        "References & Resources"
    ]
    
    for i, item in enumerate(agenda_items):
        p = content_text.paragraphs[i] if i < len(content_text.paragraphs) else content_text.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_DARK
    
    # Slide 3: Introduction
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Introduction to EV Charging Protocols"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    # Why Communication Protocols Matter
    p1 = content_text.paragraphs[0]
    p1.text = "Why Communication Protocols Matter:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    reasons = [
        "Interoperability: Seamless charging across different networks",
        "Security: Secure data exchange between EV and charging station", 
        "User Experience: Plug-and-charge convenience",
        "Grid Integration: Smart charging and demand response",
        "Billing: Automated payment processing"
    ]
    
    for i, reason in enumerate(reasons):
        p = content_text.add_paragraph()
        p.text = f"• {reason}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Evolution Timeline
    p_timeline = content_text.add_paragraph()
    p_timeline.text = "Evolution Timeline:"
    p_timeline.font.bold = True
    p_timeline.font.size = Pt(18)
    p_timeline.font.color.rgb = PRIMARY_BLUE
    
    timeline_items = [
        "2012: OCPP 1.5 released",
        "2015: OCPP 1.6 released", 
        "2015: ISO 15118-2 published",
        "2022: ISO 15118-20 (V2G) published"
    ]
    
    for item in timeline_items:
        p = content_text.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 4: OCPP 1.6 Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "OCPP 1.6 Overview"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    # What is OCPP 1.6
    p1 = content_text.paragraphs[0]
    p1.text = "What is OCPP 1.6?"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    ocpp_features = [
        "Open Charge Point Protocol 1.6",
        "WebSocket-based communication",
        "JSON message format",
        "Central System to Charge Point communication"
    ]
    
    for feature in ocpp_features:
        p = content_text.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Key Features
    p_features = content_text.add_paragraph()
    p_features.text = "Key Features:"
    p_features.font.bold = True
    p_features.font.size = Pt(18)
    p_features.font.color.rgb = SECONDARY_GREEN
    
    features = [
        "Real-time communication",
        "Transaction management",
        "Remote configuration",
        "Firmware updates",
        "Security (TLS encryption)",
        "Extensible message structure"
    ]
    
    for feature in features:
        p = content_text.add_paragraph()
        p.text = f"✅ {feature}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Limitations
    p_limitations = content_text.add_paragraph()
    p_limitations.text = "Limitations:"
    p_limitations.font.bold = True
    p_limitations.font.size = Pt(18)
    p_limitations.font.color.rgb = ACCENT_ORANGE
    
    limitations = [
        "No direct EV-to-charging station communication",
        "Manual authentication required",
        "Limited smart charging capabilities",
        "No vehicle-to-grid (V2G) support"
    ]
    
    for limitation in limitations:
        p = content_text.add_paragraph()
        p.text = f"❌ {limitation}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 5: OCPP 1.6 System Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "OCPP 1.6 System Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "OCPP 1.6 Communication Flow:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    # Create architecture diagram using text
    architecture_text = """
    ┌─────────────────┐    WebSocket    ┌──────────────────┐    HTTP/REST    ┌─────────────────┐
    │   Electric      │ ◄──────────────► │   Charging       │ ◄──────────────► │   Central       │
    │   Vehicle       │                 │   Station        │                 │   System        │
    │                 │                 │                  │                 │                 │
    │ • No direct     │                 │ • OCPP Client    │                 │ • Backend       │
    │   communication │                 │ • WebSocket      │                 │   Services      │
    │ • Manual auth   │                 │   Connection     │                 │ • Database      │
    │ • RFID/App      │                 │ • JSON Messages  │                 │ • Billing       │
    └─────────────────┘                 └──────────────────┘                 └─────────────────┘
    """
    
    p_arch = content_text.add_paragraph()
    p_arch.text = architecture_text
    p_arch.font.size = Pt(12)
    p_arch.font.color.rgb = TEXT_DARK
    p_arch.font.name = 'Courier New'
    
    # Key Components
    p_components = content_text.add_paragraph()
    p_components.text = "Key Components:"
    p_components.font.bold = True
    p_components.font.size = Pt(18)
    p_components.font.color.rgb = SECONDARY_GREEN
    
    components = [
        "Charging Station: OCPP client with WebSocket connection",
        "Central System: Backend server managing multiple stations",
        "Communication: JSON messages over WebSocket protocol",
        "Authentication: Manual via RFID, mobile app, or manual input"
    ]
    
    for component in components:
        p = content_text.add_paragraph()
        p.text = f"• {component}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 6: OCPP 1.6 Workflow
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "OCPP 1.6 Charging Workflow"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "OCPP 1.6 Charging Process:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    # Create workflow diagram using text
    workflow_text = """
    1. User arrives at charging station
    2. Manual authentication required:
       ├─ RFID card swipe
       ├─ Mobile app scan
       └─ Manual input (PIN/card)
    3. Station sends authentication to Central System
    4. Central System validates and authorizes
    5. Station starts charging session
    6. Real-time status updates via WebSocket
    7. Charging session monitoring
    8. Session termination
    9. Billing and payment processing
    10. Receipt generation
    """
    
    p_workflow = content_text.add_paragraph()
    p_workflow.text = workflow_text
    p_workflow.font.size = Pt(14)
    p_workflow.font.color.rgb = TEXT_DARK
    p_workflow.font.name = 'Courier New'
    
    # Slide 7: ISO 15118 Overview
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "ISO 15118 Overview"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    # What is ISO 15118
    p1 = content_text.paragraphs[0]
    p1.text = "What is ISO 15118?"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    iso_features = [
        "International standard for EV charging communication",
        "Direct communication between EV and charging station",
        "PLC (Power Line Communication) or wireless",
        "Plug-and-Charge (PnC) capability"
    ]
    
    for feature in iso_features:
        p = content_text.add_paragraph()
        p.text = f"• {feature}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Key Components
    p_components = content_text.add_paragraph()
    p_components.text = "Key Components:"
    p_components.font.bold = True
    p_components.font.size = Pt(18)
    p_components.font.color.rgb = SECONDARY_GREEN
    
    components = [
        "ISO 15118-1: General information and use-case definition",
        "ISO 15118-2: Network and application protocol requirements",
        "ISO 15118-3: Physical and data link layer requirements",
        "ISO 15118-4: Network and application protocol conformance testing",
        "ISO 15118-5: Physical layer and data link layer conformance testing",
        "ISO 15118-8: Physical layer and data link layer requirements for wireless communication",
        "ISO 15118-20: 2nd generation network and application protocol requirements"
    ]
    
    for component in components:
        p = content_text.add_paragraph()
        p.text = f"• {component}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 8: ISO 15118 System Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "ISO 15118 System Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "ISO 15118 Communication Flow:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    # Create architecture diagram using text
    architecture_text = """
    ┌─────────────────┐    ISO 15118    ┌──────────────────┐    OCPP/HTTP    ┌─────────────────┐
    │   Electric      │ ◄──────────────► │   Charging       │ ◄──────────────► │   Central       │
    │   Vehicle       │                 │   Station        │                 │   System        │
    │                 │                 │                  │                 │                 │
    │ • ISO 15118     │                 │ • ISO 15118      │                 │ • Backend       │
    │   Client        │                 │   Server         │                 │   Services      │
    │ • Digital       │                 │ • PLC/Wireless   │                 │ • PKI           │
    │   Certificate   │                 │ • OCPP Client    │                 │ • V2G Services  │
    └─────────────────┘                 └──────────────────┘                 └─────────────────┘
    """
    
    p_arch = content_text.add_paragraph()
    p_arch.text = architecture_text
    p_arch.font.size = Pt(12)
    p_arch.font.color.rgb = TEXT_DARK
    p_arch.font.name = 'Courier New'
    
    # Key Components
    p_components = content_text.add_paragraph()
    p_components.text = "Key Components:"
    p_components.font.bold = True
    p_components.font.size = Pt(18)
    p_components.font.color.rgb = SECONDARY_GREEN
    
    components = [
        "Electric Vehicle: ISO 15118 client with digital certificate",
        "Charging Station: ISO 15118 server + OCPP client",
        "Communication: ISO 15118 over PLC/Wireless + OCPP over WebSocket",
        "Authentication: Automatic via PKI digital certificates"
    ]
    
    for component in components:
        p = content_text.add_paragraph()
        p.text = f"• {component}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 9: ISO 15118 Workflow
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "ISO 15118 Plug-and-Charge Workflow"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "ISO 15118 Plug-and-Charge Process:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    # Create workflow diagram using text
    workflow_text = """
    1. User plugs in EV to charging station
    2. Automatic ISO 15118 handshake:
       ├─ Vehicle identification
       ├─ Digital certificate exchange
       └─ Authentication validation
    3. Station requests charging parameters from vehicle
    4. Vehicle provides:
       ├─ Battery capacity and state
       ├─ Charging preferences
       └─ Payment information
    5. Station starts charging session
    6. Real-time communication during charging
    7. V2G capabilities (if supported)
    8. Session termination
    9. Automatic billing and payment
    10. Receipt sent to vehicle/user
    """
    
    p_workflow = content_text.add_paragraph()
    p_workflow.text = workflow_text
    p_workflow.font.size = Pt(14)
    p_workflow.font.color.rgb = TEXT_DARK
    p_workflow.font.name = 'Courier New'
    
    # Slide 10: System Architecture Comparison
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "System Architecture Comparison"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "OCPP 1.6 vs ISO 15118 Architecture:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    # Create comparison diagram using text
    comparison_text = """
    OCPP 1.6 Architecture:
    ┌─────────┐    WebSocket    ┌─────────┐    HTTP/REST    ┌─────────┐
    │   EV    │ ◄──────────────► │ Station │ ◄──────────────► │ Central │
    │ (No Dir)│                 │         │                 │ System  │
    └─────────┘                 └─────────┘                 └─────────┘
    
    ISO 15118 Architecture:
    ┌─────────┐    ISO 15118    ┌─────────┐    OCPP/HTTP    ┌─────────┐
    │   EV    │ ◄──────────────► │ Station │ ◄──────────────► │ Central │
    │ (Direct)│                 │         │                 │ System  │
    └─────────┘                 └─────────┘                 └─────────┘
    """
    
    p_comp = content_text.add_paragraph()
    p_comp.text = comparison_text
    p_comp.font.size = Pt(12)
    p_comp.font.color.rgb = TEXT_DARK
    p_comp.font.name = 'Courier New'
    
    # Key Differences
    p_diffs = content_text.add_paragraph()
    p_diffs.text = "Key Architectural Differences:"
    p_diffs.font.bold = True
    p_diffs.font.size = Pt(18)
    p_diffs.font.color.rgb = SECONDARY_GREEN
    
    differences = [
        "OCPP 1.6: No direct EV communication",
        "ISO 15118: Direct EV-to-station communication",
        "OCPP 1.6: Manual authentication required",
        "ISO 15118: Automatic authentication via PKI",
        "OCPP 1.6: Limited vehicle data",
        "ISO 15118: Rich vehicle information exchange"
    ]
    
    for diff in differences:
        p = content_text.add_paragraph()
        p.text = f"• {diff}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 11: Workflow Comparison
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Workflow Comparison"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "OCPP 1.6 vs ISO 15118 Workflow:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    # Create workflow comparison using text
    workflow_comp_text = """
    OCPP 1.6 Workflow:
    1. User arrives → 2. Manual auth → 3. Station auth → 4. Start charging
    5. Monitor → 6. Stop → 7. Manual billing → 8. Receipt
    
    ISO 15118 Workflow:
    1. Plug in → 2. Auto handshake → 3. Get params → 4. Start charging
    5. Monitor → 6. Stop → 7. Auto billing → 8. Auto receipt
    """
    
    p_workflow_comp = content_text.add_paragraph()
    p_workflow_comp.text = workflow_comp_text
    p_workflow_comp.font.size = Pt(14)
    p_workflow_comp.font.color.rgb = TEXT_DARK
    p_workflow_comp.font.name = 'Courier New'
    
    # Key Differences
    p_diffs = content_text.add_paragraph()
    p_diffs.text = "Workflow Differences:"
    p_diffs.font.bold = True
    p_diffs.font.size = Pt(18)
    p_diffs.font.color.rgb = SECONDARY_GREEN
    
    workflow_diffs = [
        "Authentication: Manual vs Automatic",
        "User Experience: Multiple steps vs Plug-and-Charge",
        "Data Exchange: Limited vs Rich vehicle information",
        "Billing: Manual processing vs Automatic",
        "V2G Support: Not available vs Full support"
    ]
    
    for diff in workflow_diffs:
        p = content_text.add_paragraph()
        p.text = f"• {diff}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 12: Key Differences: Communication Architecture
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Key Differences: Communication Architecture"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "OCPP 1.6 Architecture:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    p2 = content_text.add_paragraph()
    p2.text = "[EV] ←→ [Charging Station] ←→ [Central System] ←→ [Backend Services]"
    p2.font.size = Pt(16)
    p2.font.color.rgb = TEXT_DARK
    
    p3 = content_text.add_paragraph()
    p3.text = "ISO 15118 Architecture:"
    p3.font.bold = True
    p3.font.size = Pt(18)
    p3.font.color.rgb = SECONDARY_GREEN
    
    p4 = content_text.add_paragraph()
    p4.text = "[EV] ←→ [Charging Station] ←→ [Central System] ←→ [Backend Services]"
    p4.font.size = Pt(16)
    p4.font.color.rgb = TEXT_DARK
    
    p5 = content_text.add_paragraph()
    p5.text = "     (Direct Communication)"
    p5.font.size = Pt(14)
    p5.font.color.rgb = ACCENT_ORANGE
    p5.font.italic = True
    
    # Key Differences Table
    p_table = content_text.add_paragraph()
    p_table.text = "Key Differences:"
    p_table.font.bold = True
    p_table.font.size = Pt(18)
    p_table.font.color.rgb = SECONDARY_GREEN
    
    differences = [
        "Communication Path: Central System ↔ Charging Station vs EV ↔ Charging Station",
        "Authentication: Manual (RFID, App) vs Automatic (Digital Certificate)",
        "Data Exchange: Limited vehicle data vs Rich vehicle information",
        "Security: TLS encryption vs PKI-based security"
    ]
    
    for diff in differences:
        p = content_text.add_paragraph()
        p.text = f"• {diff}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 13: Advantages for Vendors
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Advantages for Vendors"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    vendor_types = [
        ("Hardware Vendors", [
            "Higher-value charging stations",
            "Advanced features and capabilities", 
            "Competitive differentiation",
            "Future-proof technology"
        ]),
        ("Software Vendors", [
            "Rich data for analytics",
            "Advanced billing systems",
            "Grid integration services",
            "V2G platform opportunities"
        ]),
        ("Service Providers", [
            "Enhanced user experience",
            "Automated operations",
            "Reduced manual intervention",
            "New revenue streams"
        ])
    ]
    
    for vendor_type, benefits in vendor_types:
        p = content_text.add_paragraph()
        p.text = f"{vendor_type}:"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = PRIMARY_BLUE
        
        for benefit in benefits:
            p_benefit = content_text.add_paragraph()
            p_benefit.text = f"✅ {benefit}"
            p_benefit.font.size = Pt(16)
            p_benefit.font.color.rgb = TEXT_DARK
            p_benefit.level = 1
    
    # Slide 14: Advantages for Grid Operators
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Advantages for Grid Operators"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    grid_benefits = [
        ("Smart Grid Integration", [
            "Demand response capabilities",
            "Load balancing",
            "Grid stability",
            "Renewable energy integration"
        ]),
        ("Vehicle-to-Grid (V2G) Benefits", [
            "Bidirectional power flow",
            "Grid stabilization",
            "Peak shaving",
            "Frequency regulation"
        ]),
        ("Data & Analytics", [
            "Real-time load monitoring",
            "Predictive analytics",
            "Grid planning insights",
            "Efficiency optimization"
        ])
    ]
    
    for benefit_type, benefits in grid_benefits:
        p = content_text.add_paragraph()
        p.text = f"{benefit_type}:"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_GREEN
        
        for benefit in benefits:
            p_benefit = content_text.add_paragraph()
            p_benefit.text = f"✅ {benefit}"
            p_benefit.font.size = Pt(16)
            p_benefit.font.color.rgb = TEXT_DARK
            p_benefit.level = 1
    
    # Slide 15: Advantages for Users
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Advantages for Users"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    user_benefits = [
        ("Enhanced User Experience", [
            "Plug-and-Charge convenience",
            "No manual authentication needed",
            "Automatic billing",
            "Seamless roaming"
        ]),
        ("Cost Benefits", [
            "Optimized charging costs",
            "Time-of-use pricing",
            "V2G revenue opportunities",
            "Reduced transaction fees"
        ]),
        ("Environmental Benefits", [
            "Smart charging for renewables",
            "Grid decarbonization support",
            "Reduced carbon footprint",
            "Sustainable energy integration"
        ])
    ]
    
    for benefit_type, benefits in user_benefits:
        p = content_text.add_paragraph()
        p.text = f"{benefit_type}:"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_GREEN
        
        for benefit in benefits:
            p_benefit = content_text.add_paragraph()
            p_benefit.text = f"✅ {benefit}"
            p_benefit.font.size = Pt(16)
            p_benefit.font.color.rgb = TEXT_DARK
            p_benefit.level = 1
    
    # Slide 16: Implementation Challenges
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Implementation Challenges"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    challenges = [
        ("Technical Challenges", [
            "Complex PKI infrastructure",
            "Hardware compatibility",
            "Software integration",
            "Testing and certification"
        ]),
        ("Business Challenges", [
            "High implementation costs",
            "Limited vendor support",
            "Market adoption time",
            "ROI uncertainty"
        ]),
        ("Regulatory Challenges", [
            "Standards compliance",
            "Regional variations",
            "Certification requirements",
            "Interoperability testing"
        ])
    ]
    
    for challenge_type, items in challenges:
        p = content_text.add_paragraph()
        p.text = f"{challenge_type}:"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = ACCENT_ORANGE
        
        for item in items:
            p_item = content_text.add_paragraph()
            p_item.text = f"🔧 {item}"
            p_item.font.size = Pt(16)
            p_item.font.color.rgb = TEXT_DARK
            p_item.level = 1
    
    # Slide 17: Future Outlook
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Future Outlook"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    outlook_items = [
        ("Market Trends", [
            "Growing EV adoption",
            "Smart grid development",
            "V2G market expansion",
            "Renewable energy integration"
        ]),
        ("Technology Evolution", [
            "ISO 15118-20 adoption",
            "Wireless communication",
            "Advanced V2G capabilities",
            "AI integration"
        ]),
        ("Timeline Projection", [
            "2024-2025: Early adoption phase",
            "2026-2028: Mass market adoption",
            "2029-2030: Full market penetration"
        ])
    ]
    
    for item_type, items in outlook_items:
        p = content_text.add_paragraph()
        p.text = f"{item_type}:"
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.color.rgb = SECONDARY_GREEN
        
        for item in items:
            p_item = content_text.add_paragraph()
            p_item.text = f"📈 {item}"
            p_item.font.size = Pt(16)
            p_item.font.color.rgb = TEXT_DARK
            p_item.level = 1
    
    # Slide 18: Conclusion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "Key Takeaways:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    takeaways = [
        "ISO 15118 is the future of EV charging",
        "Significant advantages over OCPP 1.6",
        "Benefits all stakeholders",
        "Investment in future technology"
    ]
    
    for takeaway in takeaways:
        p = content_text.add_paragraph()
        p.text = f"🎯 {takeaway}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    p_next = content_text.add_paragraph()
    p_next.text = "Next Steps:"
    p_next.font.bold = True
    p_next.font.size = Pt(18)
    p_next.font.color.rgb = SECONDARY_GREEN
    
    next_steps = [
        "Assess current infrastructure",
        "Plan migration strategy",
        "Engage with vendors",
        "Monitor market developments"
    ]
    
    for step in next_steps:
        p = content_text.add_paragraph()
        p.text = f"➡️ {step}"
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 19: References
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = "References & Resources"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    content = slide.placeholders[1]
    content_text = content.text_frame
    
    p1 = content_text.paragraphs[0]
    p1.text = "Official Standards:"
    p1.font.bold = True
    p1.font.size = Pt(18)
    p1.font.color.rgb = PRIMARY_BLUE
    
    standards = [
        "ISO 15118-1:2019: General information and use-case definition",
        "ISO 15118-2:2016: Network and application protocol requirements",
        "ISO 15118-3:2016: Physical and data link layer requirements",
        "ISO 15118-20:2022: 2nd generation network and application protocol requirements"
    ]
    
    for standard in standards:
        p = content_text.add_paragraph()
        p.text = f"• {standard}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    p_orgs = content_text.add_paragraph()
    p_orgs.text = "Industry Organizations:"
    p_orgs.font.bold = True
    p_orgs.font.size = Pt(18)
    p_orgs.font.color.rgb = SECONDARY_GREEN
    
    orgs = [
        "CharIN e.V.: Charging Interface Initiative",
        "OCA: Open Charge Alliance",
        "IEC: International Electrotechnical Commission",
        "SAE International: Society of Automotive Engineers"
    ]
    
    for org in orgs:
        p = content_text.add_paragraph()
        p.text = f"• {org}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.level = 1
    
    # Slide 20: Thank You
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Thank You!"
    title.text_frame.paragraphs[0].font.color.rgb = PRIMARY_BLUE
    
    subtitle.text = "Questions & Discussion\n\nContact Information:\n📧 Email: [Your Email]\n📱 Phone: [Your Phone]\n🌐 Website: [Your Website]"
    subtitle.text_frame.paragraphs[0].font.color.rgb = TEXT_DARK
    
    # Save the presentation
    prs.save('ISO_15118_vs_OCPP_1.6_Presentation.pptx')
    print("✅ Presentation created successfully: ISO_15118_vs_OCPP_1.6_Presentation.pptx")

if __name__ == "__main__":
    try:
        create_iso_15118_presentation()
    except Exception as e:
        print(f"❌ Error creating presentation: {e}")
        print("Please install required packages: pip install python-pptx") 